"""
report_builder.py
Executive PowerPoint report builder for the SOT vs. cashier staffing simulation.

Design system (colors, type scale, header/footer, card/panel components) comes
from Executive_Report_Template/report_template.py, shared across this
portfolio's executive reports.

Usage from the notebook:
    from report_builder import build_executive_report
    report_path = build_executive_report(
        scenario_params, ceiling_df, sweep_df, capacity_df, breakdown_df, event_log,
        workforce_full_df, sizing_df,
        stress_lambda, sla_seconds, kitchen_sizing_target,
        capacity_ceiling_path, p99_path, breakdown_path, hourly_wait_path,
        workforce_path, sizing_path,
        output_dir=CONFIG["RESULT_DIR"],
    )
"""

import os
import sys

_TEMPLATE_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "..", "..", "Executive_Report_Template")
sys.path.insert(0, _TEMPLATE_DIR)
from report_template import (
    NAVY, GOLD, BG_ALT, LABEL_CLR, BODY_CLR, SLIDE_W, MARGIN,
    new_presentation, new_slide as _template_new_slide, add_text, add_stat_card as _template_add_stat_card,
    add_eyebrow, add_panel, add_table, add_footnote, save_report,
)
from pptx.util import Inches

EMPHASIS_CLR, DESC_CLR = "1F2933", "44515E"

KICKER = "EXECUTIVE REPORT   |   SOT vs. CASHIER STAFFING"
FOOTER = "Discrete-Event Simulation (SimPy)  |  Order + Kitchen Tandem Queue, P99 SLA"


def new_slide(prs, title_text, page_num, total_pages):
    return _template_new_slide(prs, KICKER, FOOTER, title_text, page_num, total_pages)


def add_stat_card(slide, left, top, value, label, **kwargs):
    """This deck's stat cards use a larger type scale (32/14) than the shared default (28/13)."""
    kwargs.setdefault("value_size", 32)
    kwargs.setdefault("label_size", 14)
    return _template_add_stat_card(slide, left, top, value, label, **kwargs)


def _ceiling_for(ceiling_df, label_substr):
    row = ceiling_df[ceiling_df["resource"].str.contains(label_substr, regex=False)].iloc[0]
    return row["max_stable_throughput"]


def _short_name(name):
    """'Self-Order Terminal (SOT)' -> 'SOT'; 'Human Cashier' -> 'Cashier'."""
    if "(" in name and ")" in name:
        return name[name.index("(") + 1: name.index(")")]
    return name.split()[-1]


def _workforce_row(workforce_df, name_substr, station_substr):
    mask = (workforce_df["configuration"].str.contains(name_substr, regex=False)
            & workforce_df["configuration"].str.contains(station_substr, regex=False))
    return workforce_df[mask].iloc[0]


def build_executive_report(scenario_params, ceiling_df, sweep_df, capacity_df, breakdown_df, event_log,
                            workforce_full_df, sizing_df,
                            stress_lambda, sla_seconds, kitchen_sizing_target,
                            capacity_ceiling_path, p99_path, breakdown_path, hourly_wait_path,
                            workforce_path, sizing_path, output_dir):
    """Build the 8-slide executive deck and save it as
    Executive_SOT_Staffing_Report_<YYYYMMDD>.pptx in output_dir. Returns the saved path."""
    total_pages = 8
    cashier_name = scenario_params["scenario"].iloc[0]
    sot_name = scenario_params["scenario"].iloc[1]

    kitchen_ceiling = _ceiling_for(ceiling_df, "Kitchen")
    cashier_ceiling = _ceiling_for(ceiling_df, cashier_name)
    sot_ceiling = _ceiling_for(ceiling_df, sot_name)

    cashier_capacity = capacity_df.loc[capacity_df["scenario"] == cashier_name, "p99_sla_capacity"].iloc[0]
    sot_capacity = capacity_df.loc[capacity_df["scenario"] == sot_name, "p99_sla_capacity"].iloc[0]

    cashier_row = breakdown_df[breakdown_df["scenario"] == cashier_name].iloc[0]
    sot_row = breakdown_df[breakdown_df["scenario"] == sot_name].iloc[0]

    prs = new_presentation()

    # --- Slide 1: Executive Summary ---
    s1 = new_slide(prs, "Executive Summary", 1, total_pages)
    cards = [
        (f"{kitchen_ceiling:.0f}/hr", "Shared kitchen ceiling"),
        (f"{cashier_ceiling:.0f}/hr", "Cashier register ceiling"),
        (f"{sot_ceiling:.0f}/hr", "SOT kiosk ceiling"),
        (f"{sla_seconds / 60:.0f} min", "P99 SLA target"),
    ]
    card_w, gap = 2.12, 0.17
    for i, (value, label) in enumerate(cards):
        left = MARGIN + i * (card_w + gap)
        add_stat_card(s1, left, 1.75, value, label)

    panel_w = 4.39
    add_eyebrow(s1, MARGIN, 3.44, panel_w, "KEY FINDING")
    add_panel(s1, MARGIN, 3.78, panel_w, 2.64, accent=NAVY)
    add_text(s1, MARGIN + 0.31, 4.06, panel_w - 0.56, 2.17, [
        (f"The shared kitchen's {kitchen_ceiling:.0f}/hr ceiling sits close to the {cashier_name}'s own "
         f"{cashier_ceiling:.0f}/hr register ceiling, and well below the {sot_name}'s {sot_ceiling:.0f}/hr "
         f"kiosk ceiling.", 16, False, BODY_CLR, False),
        ("The kitchen -- not either front end -- is the real constraint for both scenarios.", 16, False, BODY_CLR, False),
    ])

    right = MARGIN + panel_w + 0.22
    add_eyebrow(s1, right, 3.44, panel_w, "HEADLINE RESULT")
    add_panel(s1, right, 3.78, panel_w, 2.64, accent=GOLD)
    add_text(s1, right + 0.31, 4.06, panel_w - 0.56, 2.17, [
        (f"P99-safe capacity: {cashier_name} ~{cashier_capacity:.1f}/hr vs. {sot_name} ~{sot_capacity:.1f}/hr.", 16, False, BODY_CLR, False),
        ("A modest, not decisive, edge -- once the kitchen is a genuine co-bottleneck, extra kiosks buy "
         "little and the SOT's slower per-transaction time becomes a small net cost.", 16, False, BODY_CLR, False),
    ])

    # --- Slide 2: Capacity Ceiling by Resource ---
    s2 = new_slide(prs, "The Kitchen Is a Shared Bottleneck", 2, total_pages)
    add_eyebrow(s2, MARGIN, 1.58, 5.5, "MAX STABLE THROUGHPUT BY RESOURCE")
    s2.shapes.add_picture(capacity_ceiling_path, Inches(MARGIN), Inches(1.92), width=Inches(5.5))

    side_left = 6.28
    add_eyebrow(s2, side_left, 1.58, 3.33, "WHY THIS MATTERS")
    add_panel(s2, side_left, 1.92, 3.22, 4.12, accent=NAVY)
    add_text(s2, side_left + 0.28, 2.11, 2.72, 3.7, [
        ("Every restaurant has 3 constraints, not 1: registers, kiosks, and the kitchen behind both.", 14, False, BODY_CLR, False),
        (f"Here the kitchen caps at {kitchen_ceiling:.0f}/hr -- below the {sot_name}'s {sot_ceiling:.0f}/hr kiosk capacity.", 14, False, BODY_CLR, False),
        ("Installing more kiosks only helps up to the kitchen's own ceiling. Past that, extra front-end "
         "capacity is wasted.", 14, False, BODY_CLR, False),
    ])
    add_footnote(s2, MARGIN, 6.17, 8.61, "Kitchen and order-taking AHT figures are illustrative assumptions, checked against published QSR benchmarks -- see the notebook for sourcing.")

    # --- Slide 3: P99 vs. Demand ---
    s3 = new_slide(prs, "P99 Time Vs. Demand", 3, total_pages)
    add_eyebrow(s3, MARGIN, 1.58, 5.5, "P99 END-TO-END TIME: ORDER TO FOOD")
    s3.shapes.add_picture(p99_path, Inches(MARGIN), Inches(1.92), width=Inches(5.5))

    add_eyebrow(s3, side_left, 1.58, 3.33, "P99-SAFE CAPACITY")
    add_stat_card(s3, side_left, 1.92, f"{cashier_capacity:.1f}/hr", f"{cashier_name}", width=3.22, accent=NAVY)
    add_stat_card(s3, side_left, 3.53, f"{sot_capacity:.1f}/hr", f"{sot_name}", width=3.22, accent=GOLD)

    add_panel(s3, side_left, 5.14, 3.22, 0.9, accent=GOLD, bg=BG_ALT)
    add_text(s3, side_left + 0.24, 5.28, 2.8, 0.62, [
        ("Below this capacity, both scenarios comfortably clear the SLA; past it, P99 time explodes for both as the shared kitchen saturates.", 13, False, EMPHASIS_CLR, False),
    ])

    # --- Slide 4: Where the Time Goes ---
    s4 = new_slide(prs, "Where the Time Goes Under Stress", 4, total_pages)
    add_eyebrow(s4, MARGIN, 1.58, 4.17, f"STAGE BREAKDOWN AT {stress_lambda} CUSTOMERS/HOUR")
    s4.shapes.add_picture(breakdown_path, Inches(MARGIN), Inches(1.92), width=Inches(4.6))

    table_left = 5.5
    add_eyebrow(s4, table_left, 1.58, 4.17, "MEAN SECONDS BY STAGE")
    headers = ["Stage", cashier_name, sot_name]
    rows = [
        ("Order queue", f"{cashier_row['order_wait']:.0f}s", f"{sot_row['order_wait']:.0f}s"),
        ("Order service", f"{cashier_row['order_service']:.0f}s", f"{sot_row['order_service']:.0f}s"),
        ("Kitchen queue", f"{cashier_row['kitchen_wait']:.0f}s", f"{sot_row['kitchen_wait']:.0f}s"),
        ("Kitchen prep", f"{cashier_row['kitchen_service']:.0f}s", f"{sot_row['kitchen_service']:.0f}s"),
        ("P99 total", f"{cashier_row['p99_total']:.0f}s", f"{sot_row['p99_total']:.0f}s"),
        ("% within SLA", f"{cashier_row['sla_pct']:.1f}%", f"{sot_row['sla_pct']:.1f}%"),
    ]
    add_table(s4, table_left, 1.92, 4.17, headers, rows, first_col_frac=0.4)
    add_footnote(s4, MARGIN, 6.17, 8.61, f"Stress test at {stress_lambda}/hr, close to the shared kitchen's theoretical ceiling. Values pooled across multiple simulation replications.")

    # --- Slide 5: One Realistic Operating Day ---
    total_customers = len(event_log)
    sla_breaches = (event_log["total_time_seconds"] > sla_seconds).sum()
    cashier_log = event_log[event_log["scenario"] == cashier_name]["total_time_seconds"]
    sot_log = event_log[event_log["scenario"] == sot_name]["total_time_seconds"]

    s5 = new_slide(prs, "One Realistic Operating Day", 5, total_pages)
    add_eyebrow(s5, MARGIN, 1.58, 5.5, "MEAN WAIT TIME BY HOUR (07:00-23:00)")
    s5.shapes.add_picture(hourly_wait_path, Inches(MARGIN), Inches(1.92), width=Inches(5.5))

    add_eyebrow(s5, side_left, 1.58, 3.33, "DAILY VALIDATION")
    add_stat_card(s5, side_left, 1.92, f"{total_customers:,}", "Customers served, both scenarios", width=3.22, accent=NAVY)
    add_stat_card(s5, side_left, 3.53, f"{sla_breaches}", "SLA breaches all day", width=3.22,
                  accent=GOLD if sla_breaches == 0 else "B0413E", value_color=NAVY if sla_breaches == 0 else "B0413E")

    add_panel(s5, side_left, 5.14, 3.22, 1.2, accent=GOLD, bg=BG_ALT)
    add_text(s5, side_left + 0.24, 5.28, 2.8, 0.95, [
        (f"Mean total time: {_short_name(cashier_name)} {cashier_log.mean():.0f}s vs. {_short_name(sot_name)} {sot_log.mean():.0f}s.", 13, False, EMPHASIS_CLR, False),
        (f"Worst case: {cashier_log.max():.0f}s vs. {sot_log.max():.0f}s.", 13, False, EMPHASIS_CLR, False),
    ])
    add_footnote(s5, MARGIN, 6.17, 8.61, "One simulated day with a lunch peak and a smaller after-office/dinner peak in demand -- both scenarios hold the SLA, but the SOT runs with less headroom.")

    # --- Slide 6: Which Staffing Mix Is Most Efficient? ---
    current_stations = int(sizing_df["kitchen_stations"].min())
    cashier3 = _workforce_row(workforce_full_df, cashier_name, f"{current_stations}-station")
    kiosk3 = _workforce_row(workforce_full_df, sot_name, f"{current_stations}-station")
    combined3 = _workforce_row(workforce_full_df, "Combined", f"{current_stations}-station")

    s6 = new_slide(prs, "Which Staffing Mix Is Most Efficient?", 6, total_pages)
    add_eyebrow(s6, MARGIN, 1.58, 5.5, "P99-SAFE CAPACITY PER STAFF MEMBER")
    s6.shapes.add_picture(workforce_path, Inches(MARGIN), Inches(1.92), width=Inches(5.5))

    add_eyebrow(s6, side_left, 1.58, 3.33, f"TODAY'S KITCHEN ({current_stations} STATIONS)")
    card_h3 = 1.3
    add_stat_card(s6, side_left, 1.92, f"{cashier3['throughput_per_staff']:.1f}/hr",
                  f"{_short_name(cashier_name)} -- {cashier3['total_staff']:.0f} staff", width=3.22, height=card_h3, accent=NAVY)
    add_stat_card(s6, side_left, 1.92 + (card_h3 + 0.14), f"{kiosk3['throughput_per_staff']:.1f}/hr",
                  f"{_short_name(sot_name)} -- {kiosk3['total_staff']:.0f} staff", width=3.22, height=card_h3, accent=GOLD)
    add_stat_card(s6, side_left, 1.92 + 2 * (card_h3 + 0.14), f"{combined3['throughput_per_staff']:.1f}/hr",
                  f"Combined -- {combined3['total_staff']:.0f} staff", width=3.22, height=card_h3, accent=NAVY)
    add_footnote(s6, MARGIN, 6.17, 8.61, "Capacity per staff = P99-safe throughput / (order-taking + kitchen headcount). All three configurations deliver nearly the same throughput today -- kiosk-only does it with the fewest people.")

    # --- Slide 7: Sizing the Kitchen for Growth ---
    min_stations = int(sizing_df.loc[sizing_df["p99_seconds"] <= sla_seconds, "kitchen_stations"].min())
    kiosk9 = _workforce_row(workforce_full_df, sot_name, "9-station")
    combined9 = _workforce_row(workforce_full_df, "Combined", "9-station")

    s7 = new_slide(prs, "Sizing the Kitchen for Growth", 7, total_pages)
    add_eyebrow(s7, MARGIN, 1.58, 5.5, f"P99 TIME vs. KITCHEN STATIONS AT {kitchen_sizing_target}/HR")
    s7.shapes.add_picture(sizing_path, Inches(MARGIN), Inches(1.92), width=Inches(5.5))

    add_eyebrow(s7, side_left, 1.58, 3.33, f"TO REACH ~{kitchen_sizing_target}/HR")
    add_stat_card(s7, side_left, 1.92, f"{min_stations}", "Kitchen stations needed", width=3.22, accent=NAVY)
    add_stat_card(s7, side_left, 3.53, f"+{min_stations - current_stations}", "More hires than today", width=3.22, accent=GOLD)

    add_panel(s7, side_left, 5.14, 3.22, 1.2, accent=GOLD, bg=BG_ALT)
    add_text(s7, side_left + 0.24, 5.28, 2.8, 0.95, [
        (f"With that same {min_stations}-station kitchen, kiosk-only still tops out near {kiosk9['p99_safe_capacity']:.0f}/hr -- "
         f"only the combined front end (~{combined9['p99_safe_capacity']:.0f}/hr) can actually use the extra capacity.", 13, False, EMPHASIS_CLR, False),
    ])
    add_footnote(s7, MARGIN, 6.17, 8.61, "Growth only pays off when the kitchen and the ordering channel are expanded together -- either alone leaves capacity on the table.")

    # --- Slide 8: Recommendations ---
    s8 = new_slide(prs, "Recommendations", 8, total_pages)
    add_panel(s8, MARGIN, 1.69, 9.0, 1.03, accent=GOLD, bg=BG_ALT)
    add_text(s8, MARGIN + 0.33, 1.86, 8.39, 0.53, [
        ("PRIORITY ACTION", 14, True, GOLD, False),
        (f"Run {_short_name(sot_name)}-only today -- it matches {_short_name(cashier_name)} and Combined on capacity with the fewest staff. Only add cashiers or kitchen capacity if pursuing a specific growth target.", 16, False, EMPHASIS_CLR, False),
    ])

    add_eyebrow(s8, MARGIN, 2.97, 8.61, "KEY TAKEAWAYS")
    takeaways = [
        ("KIOSK-ONLY WINS TODAY", f"Cashier, {_short_name(sot_name)}-only, and Combined all cap out at nearly the same throughput -- {_short_name(sot_name)}-only delivers it with zero order-taking staff.", NAVY),
        ("KITCHEN FIRST, IF GROWING", "Any capacity investment should go to the kitchen before the front end -- it's the actual constraint, and combining channels alone raised capacity by ~0%.", GOLD),
        ("PAIR EXPANSION CORRECTLY", f"A {min_stations}-station kitchen unlocks ~{kitchen_sizing_target}/hr, but only when paired with the combined front end -- {_short_name(sot_name)}-only alone can't use that capacity.", NAVY),
        ("WATCH THE TAIL", "A P99 SLA is far more demanding than an average-wait target -- size every capacity decision against the tail, not the mean.", GOLD),
    ]
    card_w2, card_h2, card_gap2 = 4.39, 1.50, 0.22
    for i, (tag, text, color) in enumerate(takeaways):
        left = MARGIN + (i % 2) * (card_w2 + card_gap2)
        top = 3.39 + (i // 2) * (card_h2 + 0.19)
        add_panel(s8, left, top, card_w2, card_h2, accent=color)
        add_text(s8, left + 0.31, top + 0.10, card_w2 - 0.6, card_h2 - 0.2, [
            (tag, 14, True, color, False),
            (text, 15, False, DESC_CLR, False),
        ])

    return save_report(prs, output_dir, "Executive_SOT_Staffing_Report")
