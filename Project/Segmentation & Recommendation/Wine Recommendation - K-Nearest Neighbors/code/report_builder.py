"""
report_builder.py
Executive PowerPoint report builder for the Wine Recommendation project --
4 slides: executive summary, recommendation quality (precision@k), a sample
recommendation walkthrough with a chemical-profile radar chart, and the
full similarity landscape (PCA).

Design system (navy/gold, header/footer band, stat cards) shares the visual
identity used across this portfolio's executive reports. Output layout
(result/figures + result/slides) mirrors the Headcount Attrition and
Temperature Forecast report builders.
"""

import os
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY, GOLD = "14324B", "C9772F"
BG_LIGHT, BG_ALT, DIVIDER = "F8FAFC", "F1F5F9", "DCE3EA"
KICKER_CLR, LABEL_CLR = "9FB6C9", "6B7C8C"
BODY_CLR, FOOTER_CLR, WHITE = "2C3945", "8A98A5", "FFFFFF"

SLIDE_W, SLIDE_H, MARGIN = 10.0, 7.5, 0.5


def _rgb(hexstr):
    return RGBColor.from_string(hexstr)


def _rect(slide, left, top, width, height, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (line, size, bold, color, italic) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        run = p.runs[0]
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = _rgb(color)
    return box


def new_slide(prs, title_text, page_num, total_pages):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SLIDE_W, 1.33, NAVY)
    _rect(slide, 0, 1.33, SLIDE_W, 0.07, GOLD)
    _rect(slide, MARGIN, 6.83, SLIDE_W - 2 * MARGIN, 0.01, DIVIDER)
    _text(slide, MARGIN, 0.31, 8.61, 0.28, [("EXECUTIVE REPORT   |   WINE RECOMMENDATION", 14, True, KICKER_CLR, False)])
    _text(slide, MARGIN, 0.64, 8.61, 0.54, [(title_text, 32, True, WHITE, False)])
    _text(slide, MARGIN, 6.94, 6.67, 0.28, [("K-Nearest Neighbors  |  Cosine Similarity on Chemical Profile", 14, False, FOOTER_CLR, False)])
    _text(slide, 8.39, 6.94, 1.11, 0.28, [(f"{page_num:02d}/{total_pages:02d}", 14, True, NAVY, False)], align=PP_ALIGN.RIGHT)
    return slide


def add_stat_card(slide, left, top, value, label, accent=GOLD, bg=BG_LIGHT, value_color=NAVY, width=2.12, height=1.44):
    _rect(slide, left, top, width, height, bg)
    _rect(slide, left, top, 0.06, height, accent)
    _text(slide, left + 0.16, top + 0.08, width - 0.28, height - 0.16,
          [(value, 28, True, value_color, False), (label, 13, False, LABEL_CLR, False)])


def add_eyebrow(slide, left, top, width, text, color=LABEL_CLR):
    _text(slide, left, top, width, 0.28, [(text, 14, True, color, False)])


def add_panel(slide, left, top, width, height, accent=NAVY, bg=BG_LIGHT):
    _rect(slide, left, top, width, height, bg)
    _rect(slide, left, top, 0.06, height, accent)


def add_table(slide, left, top, width, headers, rows, header_h=0.5, row_h=0.42):
    """A simple grid styled to match the report's navy/gold design system."""
    n_cols = len(headers)
    col_w = width / n_cols

    x = left
    for label in headers:
        _rect(slide, x, top, col_w, header_h, NAVY)
        _text(slide, x + 0.12, top, col_w - 0.2, header_h, [(label, 13, True, WHITE, False)])
        x += col_w

    y = top + header_h
    for i, row_vals in enumerate(rows):
        bg = BG_LIGHT if i % 2 == 0 else BG_ALT
        x = left
        for val in row_vals:
            _rect(slide, x, y, col_w, row_h, bg)
            _text(slide, x + 0.12, y, col_w - 0.2, row_h, [(str(val), 13, False, BODY_CLR, False)])
            x += col_w
        y += row_h

    return y


def build_executive_report(
    n_wines, n_features, n_cultivars, avg_precision_at5,
    precision_curve_img, radar_img, pca_img,
    reference_label, recommended_rows,
    output_dir,
):
    """
    n_wines, n_features, n_cultivars : dataset headline counts.
    avg_precision_at5                : mean precision@5 across the whole catalog.
    precision_curve_img              : path to the precision@k chart.
    radar_img                        : path to the reference-vs-recommended radar chart.
    pca_img                          : path to the PCA similarity-landscape scatter.
    reference_label                  : the sample wine used for the walkthrough (e.g. "Wine #12 (Cultivar 0)").
    recommended_rows                 : list of (label, cultivar, similarity_str) tuples for the table.
    output_dir                       : the project's "result" directory -- the .pptx is saved to
                                        output_dir/slides, figures already live in output_dir/figures.

    Returns the saved report path.
    """
    total_pages = 4
    prs = Presentation()

    # --- Slide 1: Executive Summary ---
    s1 = new_slide(prs, "Executive Summary", 1, total_pages)
    cards = [
        (f"{n_wines}", "Wines in catalog"),
        (f"{n_features}", "Chemical attributes"),
        (f"{n_cultivars}", "Cultivars"),
        (f"{avg_precision_at5:.0%}", "Avg precision@5"),
    ]
    card_w, gap = 2.12, 0.17
    for i, (value, label) in enumerate(cards):
        left = MARGIN + i * (card_w + gap)
        add_stat_card(s1, left, 1.75, value, label)

    add_eyebrow(s1, MARGIN, 3.6, 9.0, "METHOD")
    add_panel(s1, MARGIN, 3.94, SLIDE_W - 2 * MARGIN, 2.5, accent=GOLD)
    _text(s1, MARGIN + 0.25, 4.16, SLIDE_W - 2 * MARGIN - 0.5, 2.1, [
        ("Every wine's 13 lab-measured chemical attributes (alcohol, phenols, color intensity, proline, etc.) are standardized, then compared pairwise with cosine similarity.", 14, False, BODY_CLR, False),
        ("Given a wine a taster likes, a K-Nearest Neighbors lookup finds the closest matches in that standardized space -- no other tasters or ratings required, unlike collaborative filtering.", 14, False, BODY_CLR, False),
        ("Precision@5 validates the approach: it checks whether a wine's nearest neighbors actually share its cultivar (its real, known style), using the dataset's ground-truth labels.", 14, False, BODY_CLR, False),
    ])

    # --- Slide 2: Recommendation Quality ---
    s2 = new_slide(prs, "Recommendation Quality", 2, total_pages)
    add_eyebrow(s2, MARGIN, 1.6, 9.0, "PRECISION@K -- DOES 'CHEMICALLY SIMILAR' MEAN 'SAME STYLE'?")
    if precision_curve_img and os.path.exists(precision_curve_img):
        s2.shapes.add_picture(precision_curve_img, Inches(MARGIN), Inches(2.0), width=Inches(9.0))

    # --- Slide 3: Example Recommendation ---
    s3 = new_slide(prs, "Sample Recommendation", 3, total_pages)
    add_eyebrow(s3, MARGIN, 1.55, 9.0, f"REFERENCE: {reference_label} -- CHEMICAL PROFILE MATCH")
    if radar_img and os.path.exists(radar_img):
        s3.shapes.add_picture(radar_img, Inches(MARGIN), Inches(1.9), width=Inches(5.3))

    table_left = MARGIN + 5.3 + 0.25
    table_w = SLIDE_W - MARGIN - table_left
    add_eyebrow(s3, table_left, 1.55, table_w, "TOP RECOMMENDATIONS")
    add_table(s3, table_left, 1.9, table_w,
              headers=["Wine", "Cultivar", "Similarity"],
              rows=recommended_rows)

    # --- Slide 4: Similarity Landscape ---
    s4 = new_slide(prs, "Similarity Landscape", 4, total_pages)
    add_eyebrow(s4, MARGIN, 1.55, 9.0, "ALL WINES IN PCA SPACE, COLORED BY CULTIVAR")
    if pca_img and os.path.exists(pca_img):
        s4.shapes.add_picture(pca_img, Inches(MARGIN + 0.75), Inches(1.9), width=Inches(7.5))

    slides_dir = os.path.join(output_dir, "slides")
    os.makedirs(slides_dir, exist_ok=True)
    date_str = datetime.now().strftime("%Y%m%d")
    report_path = os.path.join(slides_dir, f"Executive_Wine_Recommendation_Report_{date_str}.pptx")
    prs.save(report_path)
    return report_path
