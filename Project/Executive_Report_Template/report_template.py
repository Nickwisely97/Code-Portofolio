"""
report_template.py
Shared executive-report design system for this portfolio's PPTX decks --
colors, layout constants, and the primitive building blocks (rectangles,
text boxes, the slide header/footer band, stat cards, eyebrows, panels,
tables, footnotes).

Every project's own report_builder.py imports this module for the shared
look, then adds only the slide content and layout specific to that
project's story (which stats, which charts, how many slides). Keeping the
design system in one place means a visual tweak here propagates to every
project's report instead of needing to be hunted down across N
near-identical copies of the same helper functions.

Usage from a project's code/report_builder.py:

    import os, sys
    _TEMPLATE_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "..", "..", "Executive_Report_Template")
    sys.path.insert(0, _TEMPLATE_DIR)
    from report_template import (
        NAVY, GOLD, BG_LIGHT, BG_ALT, DIVIDER, KICKER_CLR, LABEL_CLR,
        BODY_CLR, FOOTER_CLR, WHITE, SLIDE_W, SLIDE_H, MARGIN,
        new_presentation, new_slide, add_text, add_stat_card, add_eyebrow,
        add_panel, add_table, add_footnote, save_report,
    )

    KICKER = "EXECUTIVE REPORT   |   <PROJECT NAME>"
    FOOTER = "<Method>  |  <one-line tagline>"

    def build_executive_report(..., output_dir):
        prs = new_presentation()
        s1 = new_slide(prs, KICKER, FOOTER, "Executive Summary", 1, total_pages)
        ...
        return save_report(prs, output_dir, "Executive_<Project>_Report")

The "../../../Executive_Report_Template" path assumes the standard project
layout: Project/<Category>/<Project Name>/code/report_builder.py -- three
levels up from code/ reaches Project/, where this folder lives as a sibling
of every category folder.
"""

import os
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- Design tokens -----------------------------------------------------
NAVY, GOLD = "14324B", "C9772F"
BG_LIGHT, BG_ALT, DIVIDER = "F8FAFC", "F1F5F9", "DCE3EA"
KICKER_CLR, LABEL_CLR = "9FB6C9", "6B7C8C"
BODY_CLR, FOOTER_CLR, WHITE = "2C3945", "8A98A5", "FFFFFF"

SLIDE_W, SLIDE_H, MARGIN = 10.0, 7.5, 0.5


def new_presentation():
    return Presentation()


def _rgb(hexstr):
    return RGBColor.from_string(hexstr)


def _rect(slide, left, top, width, height, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT):
    """runs : list of (line, size, bold, color_hex, italic) tuples, one per paragraph."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (line, size, bold, color, italic) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        run = p.runs[0]
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = _rgb(color)
    return box


def new_slide(prs, kicker, footer, title_text, page_num, total_pages=None):
    """
    The shared navy/gold header band + footer divider every slide starts
    from.

    kicker, footer : project-specific strings, e.g.
                      kicker = "EXECUTIVE REPORT   |   WINE RECOMMENDATION"
                      footer = "K-Nearest Neighbors  |  Cosine Similarity on Chemical Profile"
    total_pages    : if given, the page number reads "03/07"; if None, just "03".
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SLIDE_W, 1.33, NAVY)
    _rect(slide, 0, 1.33, SLIDE_W, 0.07, GOLD)
    _rect(slide, MARGIN, 6.83, SLIDE_W - 2 * MARGIN, 0.01, DIVIDER)
    add_text(slide, MARGIN, 0.31, 8.61, 0.28, [(kicker, 14, True, KICKER_CLR, False)])
    add_text(slide, MARGIN, 0.64, 8.61, 0.54, [(title_text, 32, True, WHITE, False)])
    add_text(slide, MARGIN, 6.94, 6.67, 0.28, [(footer, 14, False, FOOTER_CLR, False)])
    page_label = f"{page_num:02d}/{total_pages:02d}" if total_pages else f"{page_num:02d}"
    add_text(slide, 8.39, 6.94, 1.11, 0.28, [(page_label, 14, True, NAVY, False)], align=PP_ALIGN.RIGHT)
    return slide


def add_stat_card(slide, left, top, value, label, accent=GOLD, bg=BG_LIGHT, value_color=NAVY,
                   width=2.12, height=1.44, value_size=28, label_size=13):
    _rect(slide, left, top, width, height, bg)
    _rect(slide, left, top, 0.06, height, accent)
    add_text(slide, left + 0.16, top + 0.08, width - 0.28, height - 0.16,
          [(value, value_size, True, value_color, False), (label, label_size, False, LABEL_CLR, False)])


def add_eyebrow(slide, left, top, width, text, color=LABEL_CLR):
    add_text(slide, left, top, width, 0.28, [(text, 14, True, color, False)])


def add_panel(slide, left, top, width, height, accent=NAVY, bg=BG_LIGHT):
    _rect(slide, left, top, width, height, bg)
    _rect(slide, left, top, 0.06, height, accent)


def add_footnote(slide, left, top, width, text, color=FOOTER_CLR):
    add_text(slide, left, top, width, 0.56, [(text, 14, False, color, True)])


def add_table(slide, left, top, width, headers, rows, header_h=0.5, row_h=0.45, first_col_frac=0.24):
    """
    A grid styled to match the design system instead of pptx's default
    table look. The first column is wider, bold, and left-aligned (an
    entity label like "City" or "Wine"); the rest are centered numeric/
    text columns. Pass first_col_frac=None for plain equal-width columns.

    headers : list of column labels.
    rows    : list of tuples, one per row, same length as headers.

    Returns the y-coordinate of the table's bottom edge (inches), so
    callers can position content below it without hardcoding a height.
    """
    n_cols = len(headers)
    if first_col_frac is not None:
        first_w = width * first_col_frac
        other_w = (width - first_w) / (n_cols - 1)
        col_widths = [first_w] + [other_w] * (n_cols - 1)
    else:
        col_widths = [width / n_cols] * n_cols

    x = left
    for i, (label, w) in enumerate(zip(headers, col_widths)):
        align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER
        pad = 0.2 if i == 0 else 0
        _rect(slide, x, top, w, header_h, NAVY)
        add_text(slide, x + pad, top, w - pad, header_h, [(label, 13, True, WHITE, False)], align=align)
        x += w

    y = top + header_h
    for i, row_vals in enumerate(rows):
        bg = BG_LIGHT if i % 2 == 0 else BG_ALT
        x = left
        for j, (w, val) in enumerate(zip(col_widths, row_vals)):
            _rect(slide, x, y, w, row_h, bg)
            align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            pad = 0.2 if j == 0 else 0
            bold = j == 0
            color = NAVY if j == 0 else BODY_CLR
            add_text(slide, x + pad, y, w - pad, row_h, [(str(val), 14, bold, color, False)], align=align)
            x += w
        y += row_h

    return y


def save_report(prs, output_dir, filename_prefix):
    """Save to output_dir/slides/<filename_prefix>_<YYYYMMDD>.pptx, creating the folder if needed. Returns the saved path."""
    slides_dir = os.path.join(output_dir, "slides")
    os.makedirs(slides_dir, exist_ok=True)
    date_str = datetime.now().strftime("%Y%m%d")
    report_path = os.path.join(slides_dir, f"{filename_prefix}_{date_str}.pptx")
    prs.save(report_path)
    return report_path

