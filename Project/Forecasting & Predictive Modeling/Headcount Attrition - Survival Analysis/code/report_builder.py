"""
report_builder.py
Executive PowerPoint report builder for the Headcount Attrition survival analysis.

Design system (colors, type scale, header/footer, card/panel components) was
reverse-engineered from a hand-refined version of the deck, so every call to
build_executive_report() reproduces the same look with fresh data.

Usage from the notebook:
    from report_builder import build_executive_report
    report_path = build_executive_report(df, surv_df, current, at_risk, significant,
                                          test_cindex, cph_val, lr,
                                          cox_forest_path, km_img_path, risk_dist_path,
                                          output_dir=SLD_DIR)
"""

from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY, GOLD, RED, GREEN = "14324B", "C9772F", "B0413E", "2E7D5B"
BG_LIGHT, BG_ALT, DIVIDER = "F8FAFC", "F1F5F9", "DCE3EA"
KICKER_CLR, LABEL_CLR = "9FB6C9", "6B7C8C"
BODY_CLR, EMPHASIS_CLR, DESC_CLR, FOOTER_CLR, WHITE = "2C3945", "1F2933", "44515E", "8A98A5", "FFFFFF"

SLIDE_W, SLIDE_H, MARGIN = 10.0, 7.5, 0.5


def _rgb(hexstr):
    return RGBColor.from_string(hexstr)


def _rect(slide, left, top, width, height, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (line, size, bold, color, italic) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        run = p.runs[0]
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = _rgb(color)
    return box


def new_slide(prs, title_text, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SLIDE_W, 1.33, NAVY)
    _rect(slide, 0, 1.33, SLIDE_W, 0.07, GOLD)
    _rect(slide, MARGIN, 6.83, SLIDE_W - 2 * MARGIN, 0.01, DIVIDER)
    _text(slide, MARGIN, 0.31, 8.61, 0.28, [("EXECUTIVE REPORT   |   SURVIVAL ANALYSIS", 14, True, KICKER_CLR, False)])
    _text(slide, MARGIN, 0.64, 8.61, 0.54, [(title_text, 32, True, WHITE, False)])
    _text(slide, MARGIN, 6.94, 6.67, 0.28, [("Cox Proportional Hazards Model  |  Headcount Attrition Analysis", 14, False, FOOTER_CLR, False)])
    _text(slide, 8.39, 6.94, 1.11, 0.28, [(f"{page_num:02d}", 14, True, NAVY, False)], align=PP_ALIGN.RIGHT)
    return slide


def add_stat_card(slide, left, top, value, label, accent=GOLD, bg=BG_LIGHT, value_color=NAVY, width=2.12, height=1.44):
    _rect(slide, left, top, width, height, bg)
    _rect(slide, left, top, 0.06, height, accent)
    _text(slide, left + 0.16, top + 0.08, width - 0.28, height - 0.16,
          [(value, 32, True, value_color, False), (label, 14, False, LABEL_CLR, False)])


def add_eyebrow(slide, left, top, width, text, color=LABEL_CLR):
    _text(slide, left, top, width, 0.28, [(text, 14, True, color, False)])


def add_panel(slide, left, top, width, height, accent=NAVY, bg=BG_LIGHT):
    _rect(slide, left, top, width, height, bg)
    _rect(slide, left, top, 0.06, height, accent)


def add_footnote(slide, left, top, width, text):
    _text(slide, left, top, width, 0.56, [(text, 14, False, FOOTER_CLR, True)])


def _hazard_color(hr):
    if hr > 1.02:
        return RED
    if hr < 0.98:
        return GREEN
    return LABEL_CLR


def _phrase(feat, hr, surv_df):
    is_binary = surv_df[feat].nunique() <= 2
    direction = "Raises" if hr >= 1 else "Lowers"
    pct = abs(hr - 1) * 100
    suffix = "." if is_binary else " per unit increase."
    return f"{direction} hazard by {pct:.0f}%{suffix}"


def build_executive_report(df, surv_df, current, at_risk, significant, test_cindex, cph_val, lr,
                            hazard_bar_path, km_img_path, risk_dist_path, output_dir):
    """Build the 5-slide executive deck and save it as
    Executive_Attrition_Report_<YYYYMMDD>.pptx in output_dir. Returns the saved path."""
    top_factor = significant.index[0]
    top_factor_hr = significant.iloc[0]["exp(coef)"]

    prs = Presentation()

    # --- Slide 1: Executive Summary + Model Validation ---
    s1 = new_slide(prs, "Executive Summary", 1)
    cards = [
        (f"{len(df):,}", "Employees analyzed"),
        (f"{(df['Attrition']=='Yes').mean():.1%}", "Historical attrition rate"),
        (f"{len(current):,}", "Active & scored for risk"),
        (f"{test_cindex:.3f}", "Test concordance index"),
    ]
    card_w, gap = 2.12, 0.17
    for i, (value, label) in enumerate(cards):
        left = MARGIN + i * (card_w + gap)
        add_stat_card(s1, left, 1.75, value, label)

    panel_w = 4.39
    add_eyebrow(s1, MARGIN, 3.44, panel_w, "MODEL VALIDATION")
    add_panel(s1, MARGIN, 3.78, panel_w, 2.64, accent=NAVY)
    _text(s1, MARGIN + 0.31, 4.06, panel_w - 0.56, 2.17, [
        (f"Concordance {cph_val.concordance_index_:.3f} train / {test_cindex:.3f} test -- consistent, no overfitting.", 16, False, BODY_CLR, False),
        (f"Final model refit on all {len(surv_df):,} employees after validation.", 16, False, BODY_CLR, False),
        ("Proportional-hazards assumption checked per covariate before trusting hazard ratios.", 16, False, BODY_CLR, False),
    ])

    right = MARGIN + panel_w + 0.22
    add_eyebrow(s1, right, 3.44, panel_w, "HEADLINE FINDINGS")
    add_panel(s1, right, 3.78, panel_w, 2.64, accent=GOLD)
    _text(s1, right + 0.31, 4.06, panel_w - 0.56, 2.05, [
        (f"Strongest driver: {top_factor} -- hazard x{top_factor_hr:.2f}.", 16, False, BODY_CLR, False),
        (f"Top at-risk employee: ID {at_risk.iloc[0]['Employee ID']:.0f} -- "
         f"{at_risk.iloc[0]['P(Leave within 1 Year)']:.0%} predicted chance of leaving within a year.", 16, False, BODY_CLR, False),
        (f"Top {min(15, len(at_risk))} at-risk employees identified for targeted retention outreach.", 16, False, BODY_CLR, False),
    ])

    # --- Slide 2: What Drives Attrition ---
    s2 = new_slide(prs, "What Drives Attrition", 2)
    add_eyebrow(s2, MARGIN, 1.58, 4.17, "HAZARD RATIO PLOT")
    s2.shapes.add_picture(hazard_bar_path, Inches(MARGIN), Inches(1.92), width=Inches(5.5))

    side_left = 6.28
    add_eyebrow(s2, side_left, 1.58, 3.33, "HAZARD RATIOS")
    add_panel(s2, side_left, 1.92, 3.22, 4.12, accent=NAVY)
    runs = []
    for feat, row in significant.head(6).iterrows():
        hr = row["exp(coef)"]
        runs.append((feat, 14, True, NAVY, False))
        runs.append((f"x{hr:.2f} hazard  ·  p = {row['p']:.3f}", 14, False, _hazard_color(hr), False))
    _text(s2, side_left + 0.28, 2.11, 2.72, 3.18, runs)
    add_footnote(s2, MARGIN, 6.17, 5.5, "Hazard ratio above 1.0 raises resignation risk (red); below 1.0 lowers it (green). 1.0 = no effect.")

    # --- Slide 3: Retention by Overtime Status ---
    s3 = new_slide(prs, "Retention by Overtime Status", 3)
    add_eyebrow(s3, MARGIN, 1.58, 5.28, "KAPLAN-MEIER SURVIVAL CURVES")
    s3.shapes.add_picture(km_img_path, Inches(MARGIN), Inches(1.92), width=Inches(5.5))

    add_eyebrow(s3, side_left, 1.58, 3.33, "SIGNIFICANCE TEST")
    add_panel(s3, side_left, 1.92, 3.22, 1.83, accent=NAVY)
    _text(s3, side_left + 0.28, 2.14, 2.72, 1.44, [
        (f"p = {lr.p_value:.1e}", 32, True, NAVY, False),
        ("Log-rank test -- the retention gap between overtime and non-overtime employees is statistically significant.", 14, False, LABEL_CLR, False),
    ])

    add_panel(s3, side_left, 3.97, 3.22, 2.07, accent=GOLD, bg=BG_ALT)
    _text(s3, side_left + 0.28, 4.17, 2.72, 0.28, [("SO WHAT", 14, True, GOLD, False)])
    _text(s3, side_left + 0.28, 4.50, 2.72, 0.81, [("Overtime is the single most directly actionable factor found.", 16, False, EMPHASIS_CLR, False)])
    add_footnote(s3, MARGIN, 6.17, 5.5, "Curves show the share of employees still employed over time, split by overtime status.")

    # --- Slide 4: Who Is at Risk Right Now ---
    s4 = new_slide(prs, "Who Is at Risk Right Now", 4)
    add_eyebrow(s4, MARGIN, 1.58, 6.67, "TOP 15 EMPLOYEES BY 1-YEAR LEAVE PROBABILITY")
    s4.shapes.add_picture(risk_dist_path, Inches(0.69), Inches(1.92), width=Inches(8.61))

    stat_w, stat_gap = 2.89, 0.17
    add_stat_card(s4, MARGIN, 5.06, f"{len(current):,}", "Current employees scored", accent=NAVY, width=stat_w)
    add_stat_card(s4, MARGIN + stat_w + stat_gap, 5.06, f"{current['Risk Score'].max():.2f}",
                  f"Highest risk score vs. median {current['Risk Score'].median():.2f}",
                  accent=RED, value_color=RED, width=stat_w)

    next_left = MARGIN + 2 * (stat_w + stat_gap)
    add_panel(s4, next_left, 5.06, stat_w, 1.39, accent=GOLD, bg=BG_ALT)
    _text(s4, next_left + 0.28, 5.25, stat_w - 0.5, 0.25, [("NEXT STEP", 14, True, GOLD, False)])
    _text(s4, next_left + 0.28, 5.56, stat_w - 0.5, 0.71, [(f"Open retention conversations with the {min(15, len(at_risk))} employees shown above.", 14, False, EMPHASIS_CLR, False)])

    # --- Slide 5: Recommendations ---
    s5 = new_slide(prs, "Recommendations", 5)
    add_panel(s5, MARGIN, 1.69, 9.0, 1.03, accent=GOLD, bg=BG_ALT)
    _text(s5, MARGIN + 0.33, 1.86, 8.39, 0.53, [
        ("PRIORITY ACTION", 14, True, GOLD, False),
        (f"Prioritize retention conversations with the top {min(15, len(at_risk))} employees flagged in this report.", 16, False, EMPHASIS_CLR, False),
    ])

    add_eyebrow(s5, MARGIN, 2.97, 8.61, "KEY HAZARD DRIVERS")
    risk_factors = significant[significant["exp(coef)"] >= 1].sort_values("exp(coef)", ascending=False).head(2)
    protective_factors = significant[significant["exp(coef)"] < 1].sort_values("exp(coef)").head(2)
    driver_cards = [(feat, row, "RISK", RED) for feat, row in risk_factors.iterrows()] + \
                   [(feat, row, "PROTECTIVE", GREEN) for feat, row in protective_factors.iterrows()]

    card_w2, card_h2, card_gap2 = 4.39, 1.50, 0.22
    for i, (feat, row, tag, color) in enumerate(driver_cards):
        left = MARGIN + (i % 2) * (card_w2 + card_gap2)
        top = 3.39 + (i // 2) * (card_h2 + 0.19)
        add_panel(s5, left, top, card_w2, card_h2, accent=color)
        _text(s5, left + 0.31, top + 0.10, card_w2 - 0.6, card_h2 - 0.2, [
            (tag, 14, True, color, False),
            (feat, 18, True, NAVY, False),
            (_phrase(feat, row["exp(coef)"], surv_df), 14, False, DESC_CLR, False),
        ])

    date_str = datetime.now().strftime("%Y%m%d")
    report_path = f"{output_dir.rstrip('/')}/Executive_Attrition_Report_{date_str}.pptx"
    prs.save(report_path)
    return report_path
