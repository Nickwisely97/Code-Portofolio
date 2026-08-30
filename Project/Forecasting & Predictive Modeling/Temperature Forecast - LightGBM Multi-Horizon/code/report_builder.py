"""
report_builder.py
Executive PowerPoint report builder for the Temperature Forecast pipeline.

Slide order is deliberately temperature-first: the executive summary and the
per-city forecast fans lead the deck since that's what a reader wants to know
first; backtest accuracy (MAE/RMSE/R² -- validation of the model itself, not
the forecast) is pushed to a single Model Performance slide at the very end.

Design system (navy/gold, header/footer band, stat cards) shares the visual
identity used across this portfolio's executive reports.
"""

import os
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY, GOLD = "14324B", "C9772F"
BG_LIGHT, BG_ALT, DIVIDER = "F8FAFC", "F1F5F9", "DCE3EA"
KICKER_CLR, LABEL_CLR = "9FB6C9", "6B7C8C"
BODY_CLR, FOOTER_CLR, WHITE = "2C3945", "8A98A5", "FFFFFF"

SLIDE_W, SLIDE_H, MARGIN = 10.0, 7.5, 0.5


def _rgb(hexstr):
    return RGBColor.from_string(hexstr)


def _rect(slide, left, top, width, height, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (line, size, bold, color, italic) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        run = p.runs[0]
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = _rgb(color)
    return box


def new_slide(prs, title_text, page_num, total_pages):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SLIDE_W, 1.33, NAVY)
    _rect(slide, 0, 1.33, SLIDE_W, 0.07, GOLD)
    _rect(slide, MARGIN, 6.83, SLIDE_W - 2 * MARGIN, 0.01, DIVIDER)
    _text(slide, MARGIN, 0.31, 8.61, 0.28, [("EXECUTIVE REPORT   |   TEMPERATURE FORECAST", 14, True, KICKER_CLR, False)])
    _text(slide, MARGIN, 0.64, 8.61, 0.54, [(title_text, 32, True, WHITE, False)])
    _text(slide, MARGIN, 6.94, 6.67, 0.28, [("LightGBM Direct Multi-Horizon Forecast  |  14-Day Outlook", 14, False, FOOTER_CLR, False)])
    _text(slide, 8.39, 6.94, 1.11, 0.28, [(f"{page_num:02d}/{total_pages:02d}", 14, True, NAVY, False)], align=PP_ALIGN.RIGHT)
    return slide


def add_stat_card(slide, left, top, value, label, accent=GOLD, bg=BG_LIGHT, value_color=NAVY, width=2.12, height=1.44):
    _rect(slide, left, top, width, height, bg)
    _rect(slide, left, top, 0.06, height, accent)
    _text(slide, left + 0.16, top + 0.08, width - 0.28, height - 0.16,
          [(value, 28, True, value_color, False), (label, 13, False, LABEL_CLR, False)])


def add_eyebrow(slide, left, top, width, text, color=LABEL_CLR):
    _text(slide, left, top, width, 0.28, [(text, 14, True, color, False)])


def add_panel(slide, left, top, width, height, accent=NAVY, bg=BG_LIGHT):
    _rect(slide, left, top, width, height, bg)
    _rect(slide, left, top, 0.06, height, accent)


def add_city_table(slide, left, top, width, headers, rows, header_h=0.5, row_h=0.45):
    """
    A City x N-column grid styled to match the report's navy/gold design
    system instead of pptx's default table look. Used for both the 14-day
    temperature outlook (slide 1) and the backtest accuracy table (last slide).

    headers : list of column labels; the first column ("City") is left-aligned,
              the rest are centered.
    rows    : list of tuples, one per city, same length as headers.
    """
    n_cols = len(headers)
    first_w = width * 0.24
    other_w = (width - first_w) / (n_cols - 1)
    col_widths = [first_w] + [other_w] * (n_cols - 1)

    x = left
    for label, w in zip(headers, col_widths):
        align = PP_ALIGN.LEFT if label == "City" else PP_ALIGN.CENTER
        pad = 0.2 if label == "City" else 0
        _rect(slide, x, top, w, header_h, NAVY)
        _text(slide, x + pad, top, w - pad, header_h, [(label, 13, True, WHITE, False)], align=align)
        x += w

    y = top + header_h
    for i, row_vals in enumerate(rows):
        bg = BG_LIGHT if i % 2 == 0 else BG_ALT
        x = left
        for label, w, val in zip(headers, col_widths, row_vals):
            _rect(slide, x, y, w, row_h, bg)
            align = PP_ALIGN.LEFT if label == "City" else PP_ALIGN.CENTER
            pad = 0.2 if label == "City" else 0
            bold = label == "City"
            color = NAVY if label == "City" else BODY_CLR
            _text(slide, x + pad, y, w - pad, row_h, [(str(val), 14, bold, color, False)], align=align)
            x += w
        y += row_h

    return y  # bottom edge of the table


def build_executive_report(df_summary, results, output_dir):
    """
    df_summary : DataFrame with columns Station, MAE (°C), RMSE (°C), MAPE (%), R².
    results    : dict {station: {"img_fan": path, "live": DataFrame, ...}} --
                 img_fan is the forecast-fan chart, live is the live 14-day
                 forecast frame (column "predicted_actual") used for the
                 outlook stats (avg/max/min).
    output_dir : the project's "result" directory -- the .pptx is saved to
                 output_dir/slides, mirroring the Headcount Attrition report's
                 figures/ + slides/ layout.

    Returns the saved report path.
    """
    stations = df_summary["Station"].tolist()
    total_pages = 1 + len(stations) + 1  # summary + one per city + performance

    best  = df_summary.loc[df_summary["MAE (°C)"].idxmin()]
    worst = df_summary.loc[df_summary["MAE (°C)"].idxmax()]
    avg_mae = df_summary["MAE (°C)"].mean()

    outlook = {}
    for station in stations:
        live = results.get(station, {}).get("live")
        if live is not None and not live.empty:
            outlook[station] = (
                live["predicted_actual"].mean(),
                live["predicted_actual"].max(),
                live["predicted_actual"].min(),
            )
        else:
            outlook[station] = (None, None, None)

    prs = Presentation()

    # --- Slide 1: Executive Summary -- 14-day temperature outlook by city ---
    # Temperature first: this is the headline question, so it leads the deck.
    s1 = new_slide(prs, "Executive Summary", 1, total_pages)
    cards = [
        (f"{len(stations)}", "Cities forecast"),
        ("14 days", "Forecast horizon"),
    ]
    card_w, gap = 2.12, 0.17
    for i, (value, label) in enumerate(cards):
        left = MARGIN + i * (card_w + gap)
        add_stat_card(s1, left, 1.75, value, label)

    add_eyebrow(s1, MARGIN, 3.44, 9.0, "14-DAY TEMPERATURE OUTLOOK BY CITY")
    outlook_rows = [
        (station,
         f"{avg_t:.1f}°C" if avg_t is not None else "n/a",
         f"{max_t:.1f}°C" if max_t is not None else "n/a",
         f"{min_t:.1f}°C" if min_t is not None else "n/a")
        for station, (avg_t, max_t, min_t) in outlook.items()
    ]
    table_bottom = add_city_table(
        s1, MARGIN, 3.78, 9.0,
        headers=["City", "Avg Forecast Temp", "Max Forecast Temp", "Min Forecast Temp"],
        rows=outlook_rows,
    )

    _text(s1, MARGIN, table_bottom + 0.18, 9.0, 0.4, [
        ("Model performance (backtest accuracy) is on the last slide of this deck.", 13, False, LABEL_CLR, True),
    ])

    # --- Slides 2..N: one per city -- live forecast fan only ---
    img_w = 9.0
    fan_top = 2.35
    img_h = img_w * 6 / 14  # matches plot_forecast_fan's figsize=(14, 6)

    for i, (_, row) in enumerate(df_summary.iterrows(), start=2):
        station = row["Station"]
        res = results.get(station, {})
        s = new_slide(prs, f"{station} -- 14-Day Forecast", i, total_pages)

        add_eyebrow(s, MARGIN, 1.7, 9.0, "LIVE 14-DAY FORECAST FAN")
        fan_path = res.get("img_fan", "")
        if fan_path and os.path.exists(fan_path):
            s.shapes.add_picture(fan_path, Inches(MARGIN), Inches(fan_top), width=Inches(img_w))

    # --- Final slide: Model Performance -- backtest accuracy by city ---
    # Pushed to the end on purpose: this validates the model, it isn't the
    # forecast itself, so it shouldn't compete with the temperature story.
    sN = new_slide(prs, "Model Performance", total_pages, total_pages)
    add_eyebrow(sN, MARGIN, 1.75, 9.0, "BACKTEST ACCURACY BY CITY (WALK-FORWARD ROLLBACK TESTING)")
    perf_rows = [
        (r["Station"], f"{r['MAE (°C)']:.2f}°C", f"{r['RMSE (°C)']:.2f}°C", f"{r['R²']:.2f}")
        for _, r in df_summary.iterrows()
    ]
    perf_bottom = add_city_table(
        sN, MARGIN, 2.1, 9.0,
        headers=["City", "MAE", "RMSE", "R-squared"],
        rows=perf_rows,
    )

    add_eyebrow(sN, MARGIN, perf_bottom + 0.25, 9.0, "HEADLINE FINDINGS")
    add_panel(sN, MARGIN, perf_bottom + 0.57, SLIDE_W - 2 * MARGIN, 1.3, accent=GOLD)
    _text(sN, MARGIN + 0.25, perf_bottom + 0.75, SLIDE_W - 2 * MARGIN - 0.5, 1.0, [
        (f"Most accurate: {best['Station']} ({best['MAE (°C)']:.2f}°C MAE). Least accurate: {worst['Station']} ({worst['MAE (°C)']:.2f}°C MAE).", 14, False, BODY_CLR, False),
        (f"Average backtest MAE across all cities: {avg_mae:.2f}°C, from many historical origins each forecast 14 days ahead and scored against what actually happened.", 14, False, BODY_CLR, False),
    ])

    slides_dir = os.path.join(output_dir, "slides")
    os.makedirs(slides_dir, exist_ok=True)
    date_str = datetime.now().strftime("%Y%m%d")
    report_path = os.path.join(slides_dir, f"Executive_Temperature_Report_{date_str}.pptx")
    prs.save(report_path)
    return report_path
