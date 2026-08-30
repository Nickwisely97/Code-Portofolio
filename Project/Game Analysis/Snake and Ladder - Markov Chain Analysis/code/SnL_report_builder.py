"""
SnL_report_builder.py
Executive PowerPoint report builder for the Snake & Ladder Markov chain
analysis -- 5 slides: executive summary, board heatmap, snake/ladder
impact, finish-time distribution, and the dice-count comparison.

Design system (navy/gold, header/footer band, stat cards, tables) comes from
Executive_Report_Template/report_template.py, shared across this portfolio's
executive reports -- this file only adds the content specific to this project.
"""

import os
import sys

_TEMPLATE_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "..", "..", "Executive_Report_Template")
sys.path.insert(0, _TEMPLATE_DIR)
from report_template import (
    GOLD, LABEL_CLR, BODY_CLR, SLIDE_W, MARGIN,
    new_presentation, new_slide, add_text, add_stat_card, add_eyebrow, add_panel, save_report,
)
from pptx.util import Inches

KICKER = "EXECUTIVE REPORT   |   SNAKE & LADDER"
FOOTER = "Absorbing Markov Chain  |  Fundamental Matrix Analysis"


def build_executive_report(
    n_dice, board_size, expected_steps_from_start, n_ladders, n_snakes,
    heatmap_img, impact_img, impact_findings,
    finish_img, finish_stats,
    dice_comparison_img,
    output_dir,
):
    """
    n_dice, board_size, expected_steps_from_start : headline numbers.
    n_ladders, n_snakes                            : board layout counts.
    heatmap_img, impact_img, finish_img,
    dice_comparison_img                             : chart paths (from result/figures).
    impact_findings                                 : list of up to 3 strings (top snake/ladder findings).
    finish_stats                                     : dict with keys "median", "p90", "mean", "std".
    output_dir                                       : the project's "result" directory -- the .pptx is
                                                         saved to output_dir/slides.

    Returns the saved report path.
    """
    total_pages = 5
    prs = new_presentation()

    # --- Slide 1: Executive Summary ---
    s1 = new_slide(prs, KICKER, FOOTER, "Executive Summary", 1, total_pages)
    cards = [
        (f"{expected_steps_from_start:.1f}", "Expected turns from square 1"),
        (f"{board_size}", "Board squares"),
        (f"{n_ladders}", "Ladders"),
        (f"{n_snakes}", "Snakes"),
    ]
    card_w, gap = 2.12, 0.17
    for i, (value, label) in enumerate(cards):
        left = MARGIN + i * (card_w + gap)
        add_stat_card(s1, left, 1.75, value, label)

    add_eyebrow(s1, MARGIN, 3.6, 9.0, "METHOD")
    add_panel(s1, MARGIN, 3.94, SLIDE_W - 2 * MARGIN, 2.5, accent=GOLD)
    add_text(s1, MARGIN + 0.25, 4.16, SLIDE_W - 2 * MARGIN - 0.5, 2.1, [
        (f"The board is modeled as an absorbing Markov chain, with the transition matrix built directly from the board layout and a {n_dice}-dice roll distribution -- not a fixed spreadsheet.", 14, False, BODY_CLR, False),
        ("The fundamental matrix N = (I - Q)^-1 gives the expected number of turns to finish from every square at once.", 14, False, BODY_CLR, False),
        ("Because the matrix is built in code, the same analysis re-runs for any dice count -- see the dice-count comparison on the final slide.", 14, False, BODY_CLR, False),
    ])

    # --- Slide 2: Board Heatmap ---
    s2 = new_slide(prs, KICKER, FOOTER, "Board Heatmap", 2, total_pages)
    add_eyebrow(s2, MARGIN, 1.55, 9.0, "EXPECTED STEPS TO FINISH, BY SQUARE")
    if heatmap_img and os.path.exists(heatmap_img):
        img_h = 4.75
        img_w = img_h * 18 / 16  # matches plot_board_heatmap's figsize=(18, 16)
        img_left = MARGIN + (SLIDE_W - 2 * MARGIN - img_w) / 2
        s2.shapes.add_picture(heatmap_img, Inches(img_left), Inches(1.9), height=Inches(img_h))

    # --- Slide 3: Snake & Ladder Impact ---
    s3 = new_slide(prs, KICKER, FOOTER, "Snake & Ladder Impact", 3, total_pages)
    add_eyebrow(s3, MARGIN, 1.55, 9.0, "HOW MANY TURNS EACH ELEMENT SAVES OR COSTS")
    if impact_img and os.path.exists(impact_img):
        img_h = 4.65
        img_w = img_h * 16 / 10  # matches plot_snl_impact's figsize=(16, 10)
        img_left = MARGIN + (SLIDE_W - 2 * MARGIN - img_w) / 2
        s3.shapes.add_picture(impact_img, Inches(img_left), Inches(1.9), height=Inches(img_h))
    if impact_findings:
        add_text(s3, MARGIN, 6.6, SLIDE_W - 2 * MARGIN, 0.22,
                 [(impact_findings[0], 12, False, LABEL_CLR, True)])

    # --- Slide 4: Finish-Time Distribution ---
    s4 = new_slide(prs, KICKER, FOOTER, "Finish-Time Distribution", 4, total_pages)
    add_eyebrow(s4, MARGIN, 1.55, 5.6, "WHEN DOES THE GAME ACTUALLY END?")
    if finish_img and os.path.exists(finish_img):
        s4.shapes.add_picture(finish_img, Inches(MARGIN), Inches(1.9), width=Inches(5.6))

    panel_left = MARGIN + 5.6 + 0.22
    panel_w = SLIDE_W - MARGIN - panel_left
    add_eyebrow(s4, panel_left, 1.55, panel_w, "PERCENTILES")
    stat_cards = [
        (f"{finish_stats['median']:.0f}", "Median (50%) turns"),
        (f"{finish_stats['p90']:.0f}", "90th percentile turns"),
        (f"{finish_stats['mean']:.1f}", "Mean turns"),
        (f"{finish_stats['std']:.1f}", "Std. deviation"),
    ]
    for i, (value, label) in enumerate(stat_cards):
        top = 1.9 + i * (1.05 + 0.14)
        add_stat_card(s4, panel_left, top, value, label, width=panel_w, height=1.05)

    # --- Slide 5: Dice Count Comparison ---
    s5 = new_slide(prs, KICKER, FOOTER, "Dice Count Comparison", 5, total_pages)
    add_eyebrow(s5, MARGIN, 1.55, 9.0, "HOW THE NUMBER OF DICE CHANGES GAME LENGTH")
    if dice_comparison_img and os.path.exists(dice_comparison_img):
        img_h = 4.6
        img_w = img_h * 10 / 6  # matches plot_dice_comparison's figsize=(10, 6)
        img_left = MARGIN + (SLIDE_W - 2 * MARGIN - img_w) / 2
        s5.shapes.add_picture(dice_comparison_img, Inches(img_left), Inches(2.0), height=Inches(img_h))

    return save_report(prs, output_dir, "Executive_SnakeLadder_Report")
